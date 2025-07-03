import os
import re
import tempfile
import requests
import pandas as pd
from urllib.parse import urlparse
from bs4 import BeautifulSoup
from pptx import Presentation
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
from utils.reader import read_file  # handles .txt, .docx etc.

# Path to Tesseract executable (update this path if installed elsewhere)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# -------------------------
# 🖼️ Image OCR Reader
# -------------------------
def read_image_file(file_path):
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        print("📸 OCR Extracted Text:\n", text)
        return text
    except Exception as e:
        print(f"❌ Error reading image: {e}")
        return ""

# -------------------------
# 📄 PDF OCR Reader
# -------------------------
def read_pdf_with_ocr(file_path):
    try:
        images = convert_from_path(file_path)
        text = ""
        for i, img in enumerate(images):
            page_text = pytesseract.image_to_string(img)
            text += f"\n--- Page {i+1} ---\n" + page_text
        print("📄 PDF OCR Text:\n", text)
        return text
    except Exception as e:
        print(f"❌ PDF OCR failed: {e}")
        return ""

# -------------------------
# 🧾 PPTX Reader
# -------------------------
def read_pptx_file(file_path):
    try:
        text = []
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"❌ Error reading PPTX: {e}")
        return ""

# -------------------------
# 🔗 Convert Google Drive / Docs link to direct download
# -------------------------
def convert_drive_url_to_direct(url):
    file_id = None
    if "docs.google.com/spreadsheets" in url:
        if "/d/" in url:
            file_id = url.split("/d/")[1].split("/")[0]
        return f"https://docs.google.com/spreadsheets/d/{file_id}/export?format=csv"
    if "/file/d/" in url:
        file_id = url.split("/file/d/")[1].split("/")[0]
    elif "open?id=" in url:
        file_id = url.split("open?id=")[1].split("&")[0]
    elif "uc?id=" in url:
        file_id = url.split("uc?id=")[1].split("&")[0]
    return f"https://drive.google.com/uc?export=download&id={file_id}" if file_id else url

# -------------------------
# 📎 Infer filename from URL/headers
# -------------------------
def get_filename_from_url_or_header(url, response):
    content_disposition = response.headers.get('Content-Disposition', '')
    if 'filename=' in content_disposition:
        return content_disposition.split('filename=')[1].strip('"')

    # Google Sheets
    if "docs.google.com/spreadsheets" in url:
        match = re.search(r"/d/([a-zA-Z0-9_-]+)", url)
        return f"sheet_{match.group(1)}.csv" if match else "sheet.csv"

    # Google Slides or Docs
    if "docs.google.com/presentation" in url or "docs.google.com/document" in url:
        try:
            soup = BeautifulSoup(response.text, "html.parser")
            title = soup.title.string.strip() if soup.title else "presentation"
            return re.sub(r'[^a-zA-Z0-9_.-]', '_', title) + ".pptx"
        except Exception as e:
            print(f"⚠️ Could not parse HTML title: {e}")
            return "google_presentation.pptx"

    # Generic HTML-based Drive preview
    if "drive.google.com" in url and "text/html" in response.headers.get("Content-Type", ""):
        try:
            soup = BeautifulSoup(response.text, "html.parser")
            title = soup.title.string.strip() if soup.title else "downloaded"
            return re.sub(r'[^a-zA-Z0-9_.-]', '_', title) + ".txt"
        except Exception as e:
            print(f"⚠️ Could not parse title from Drive: {e}")
            return "drive_file.txt"

    # Fallback
    parsed = urlparse(url)
    filename = os.path.basename(parsed.path)
    return filename or "downloaded.txt"

# -------------------------
# 🔻 Download file to temp path
# -------------------------
def download_file_from_url(url):
    try:
        if "drive.google.com" in url or "docs.google.com" in url:
            url = convert_drive_url_to_direct(url)
        response = requests.get(url, stream=True, allow_redirects=True, timeout=60)
        response.raise_for_status()
        filename = get_filename_from_url_or_header(url, response)
        ext = os.path.splitext(filename)[1] or ".txt"
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
            for chunk in response.iter_content(chunk_size=8192):
                tmp_file.write(chunk)
            return tmp_file.name, filename
    except Exception as e:
        print(f"❌ Error downloading: {e}")
        return None, None

# -------------------------
# 📊 Read CSV or Excel
# -------------------------
def read_csv_excel_file(file_path, filename):
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".csv" or "sheet_" in filename:
            for encoding in ['utf-8', 'latin-1', 'cp1252', 'utf-16']:
                for sep in [',', ';', '\t', '|']:
                    try:
                        df = pd.read_csv(file_path, encoding=encoding, sep=sep)
                        if df.shape[1] > 1 and not df.empty:
                            return clean_dataframe(df)
                    except:
                        continue
            try:
                df = pd.read_csv(file_path, engine='python')
                return clean_dataframe(df)
            except Exception as e:
                print(f"❌ CSV fallback failed: {e}")
                return None
        elif ext in [".xlsx", ".xls"]:
            try:
                df = pd.read_excel(file_path, engine='openpyxl' if ext == '.xlsx' else 'xlrd')
                return clean_dataframe(df)
            except Exception as e:
                print(f"❌ Excel read failed: {e}")
                return None
        return None
    except Exception as e:
        print(f"❌ Error reading structured file: {e}")
        return None

# -------------------------
# 🧹 DataFrame cleaner
# -------------------------
def clean_dataframe(df):
    df = df.dropna(how='all').loc[:, ~df.isnull().all()]
    for col in df.columns:
        if 'Unnamed' in str(col) and df[col].isna().mean() > 0.95:
            df.drop(col, axis=1, inplace=True)
    df.columns = [str(col).strip() for col in df.columns]
    if len(df) > 50000:
        df = df.head(50000)
    print(f"📊 Cleaned DataFrame: {df.shape}")
    return df

# -------------------------
# 📚 Main universal loader
# -------------------------
def read_file_text(path_or_url):
    is_url = path_or_url.startswith("http")
    temp_path = None

    try:
        if is_url:
            temp_path, filename = download_file_from_url(path_or_url)
            if not temp_path:
                return None, "unknown.txt", "unknown"
            path = temp_path
        else:
            filename = os.path.basename(path_or_url)
            path = path_or_url

        ext = os.path.splitext(filename)[1].lower()
        print(f"📁 Processing: {filename} (ext: {ext})")

        if ext in [".csv", ".xlsx", ".xls"] or "sheet_" in filename:
            df = read_csv_excel_file(path, filename)
            return (df, filename, "structured") if df is not None else (None, filename, "structured")

        elif ext == ".pptx":
            try:
                ppt_text = read_pptx_file(path)
                if ppt_text.strip():
                    return ppt_text, filename, "text"
                else:
                    raise Exception("Empty PPTX content.")
            except Exception as pptx_error:
                print(f"⚠️ Falling back to HTML parsing for PPTX: {pptx_error}")
                try:
                    with open(path, "r", encoding="utf-8", errors="ignore") as f:
                        html = f.read()
                    soup = BeautifulSoup(html, "html.parser")
                    text = soup.get_text(separator="\n")
                    return text.strip(), filename, "text"
                except Exception as html_error:
                    print(f"❌ Failed to extract text from HTML fallback: {html_error}")
                    return None, filename, "text"

        elif ext in [".png", ".jpg", ".jpeg"]:
            text = read_image_file(path)
            return (text, filename, "text") if text.strip() else (None, filename, "text")

        elif ext == ".pdf":
            text = read_pdf_with_ocr(path)
            return (text, filename, "text") if text.strip() else (None, filename, "text")

        else:
            content = read_file(path)
            return (content, filename, "text") if content.strip() else (None, filename, "text")

    except Exception as e:
        print(f"❌ Error processing file: {e}")
        return None, "error.txt", "unknown"
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except:
                pass
